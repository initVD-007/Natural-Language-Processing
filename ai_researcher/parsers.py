import pypdf
from pptx import Presentation
import google.generativeai as genai
import os
import re
from google.api_core.exceptions import GoogleAPIError
from tenacity import retry, wait_exponential, stop_after_attempt, retry_if_exception_type

class AgentA_Parser:
    def __init__(self, api_key: str):
        if not api_key:
            raise ValueError("Google API Key is required.")
        genai.configure(api_key=api_key)
        self.model = genai.GenerativeModel('gemini-3-flash-preview')

    def extract_text(self, uploaded_file):
        """Extracts text from PDF or PPTX bytes."""
        extension = uploaded_file.name.split('.')[-1].lower()
        text = ""
        
        try:
            if extension == "pdf":
                reader = pypdf.PdfReader(uploaded_file)
                for page in reader.pages:
                    text += page.extract_text() or ""
            elif extension in ["pptx", "ppt"]:
                prs = Presentation(uploaded_file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + " "
        except Exception as e:
            raise Exception(f"Failed to parse document: {str(e)}")
            
        return text

    @retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=2, min=4, max=15), retry=retry_if_exception_type(GoogleAPIError))
    def analyze_content(self, raw_text):
        """Uses LLM to deeply analyze document structure and generate research queries."""
        
        # Safely truncate if too large to prevent token explosion
        safe_text = raw_text[:30000]
        
        prompt = f"""
You are Agent A (The Intelligent Parser).

Your role is to deeply analyze the provided educational document.
The content may belong to any domain (e.g., Mathematics, Finance, Artificial Intelligence, Computer Science, etc.).

Perform the following tasks carefully and concisely:

1. MAIN SUBJECT
   - Identify the single primary subject of the document.
   - Provide a one-line clear definition of that subject.

2. CORE TOPICS
   - Extract the top 5 most important core topics discussed.
   - For each topic:
       • Provide a short explanation (1–2 sentences).
       • Mention why it is important in the context of the document.

3. 2026 UPDATE SEARCH QUERIES
   - Generate exactly 3 highly specific and research-focused search queries.
   - Each query must:
       • Include the year "2026"
       • Be designed to find recent advancements, updates, or new research
       • Be clear, targeted, and suitable for academic or technical search engines

CRITICAL: Output Format (strictly follow this markdown structure without adding extra text or commentary):

MAIN SUBJECT:
<subject name>
<one-line definition>

CORE TOPICS:
1. <Topic Name>
   - Explanation: <Explanation text>
   - Importance: <Importance text>

2. <Topic Name>
   - Explanation: <Explanation text>
   - Importance: <Importance text>

3. <Topic Name>
   - Explanation: <Explanation text>
   - Importance: <Importance text>

4. <Topic Name>
   - Explanation: <Explanation text>
   - Importance: <Importance text>

5. <Topic Name>
   - Explanation: <Explanation text>
   - Importance: <Importance text>

SEARCH QUERIES FOR 2026 UPDATES:
1. "<query 1>"
2. "<query 2>"
3. "<query 3>"

DOCUMENT TEXT:
{safe_text}
"""
        try:
            response = self.model.generate_content(prompt)
            return response.text
        except Exception as e:
            raise Exception(f"LLM Analysis failed: {str(e)}")

    def get_search_queries(self, report_text):
        """Extracts the 3 queries from the LLM report using Regex."""
        queries = re.findall(r'"([^"]*2026[^"]*)"', report_text)

        # return only first 3 found
        return queries[:3]
